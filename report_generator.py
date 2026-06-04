# =============================================================================
# report_generator.py — Excel, PDF, and PPTX Report Generator (v2.0)
# =============================================================================

import os
import re
import warnings
from datetime import datetime

import jinja2
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from bs4 import XMLParsedAsHTMLWarning

warnings.filterwarnings("ignore", category=XMLParsedAsHTMLWarning)

OUTPUT_DIR = "output"

# Brand colours
NAVY  = RGBColor(0x1a, 0x1a, 0x2e)
RED   = RGBColor(0xE2, 0x4B, 0x4A)
AMBER = RGBColor(0xEF, 0x9F, 0x27)
GREEN = RGBColor(0x3B, 0x6D, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0xF5, 0xF5, 0xF5)


def ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def sanitize(text):
    if not isinstance(text, str):
        text = str(text)
    return text.encode("ascii", "replace").decode("ascii").replace("?", " ").strip()


# ─────────────────────────────────────────────────────────────
#  EXCEL
# ─────────────────────────────────────────────────────────────

def generate_excel(ticker, analysis, comparison, latest_date, previous_date):
    ensure_output_dir()

    wb = Workbook()
    findings = analysis["findings"]
    summary  = analysis["summary"]

    high_count   = summary["by_severity"].get("HIGH", 0)
    medium_count = summary["by_severity"].get("MEDIUM", 0)
    low_count    = summary["by_severity"].get("LOW", 0)
    total        = summary["total"]
    new_kws      = comparison.get("new_keywords", [])
    avg_sent     = summary.get("avg_sentiment", 0.0)

    nav_fill  = PatternFill(start_color="1a1a2e", end_color="1a1a2e", fill_type="solid")
    nav_font  = Font(color="FFFFFF", bold=True, size=10)
    thin      = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"),  bottom=Side(style="thin"),
    )

    def _hdr_cell(ws, row, col, value):
        c = ws.cell(row=row, column=col, value=value)
        c.fill = nav_fill; c.font = nav_font
        c.border = thin
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        return c

    def _data_cell(ws, row, col, value, fill=None):
        c = ws.cell(row=row, column=col, value=value)
        c.border = thin
        c.alignment = Alignment(wrap_text=True, vertical="top")
        if fill:
            c.fill = fill
        return c

    # ── Sheet 1: Quick Brief ──────────────────────────────────
    ws1 = wb.active
    ws1.title = "Quick Brief"

    # Rows 1-4: company info block
    meta = [
        ("Company / Ticker", ticker),
        ("Exchange / Source", comparison.get("exchange", "Unknown")),
        ("Filing Date", str(latest_date or "N/A")),
        ("Previous Filing", str(previous_date or "N/A")),
        ("Verdict",
         "CRITICAL RISK" if high_count > 0 else
         "ELEVATED RISK" if medium_count > 3 else "LOW RISK"),
    ]
    for r, (k, v) in enumerate(meta, start=1):
        c1 = ws1.cell(row=r, column=1, value=k)
        c1.font = Font(bold=True, color="FFFFFF"); c1.fill = nav_fill; c1.border = thin
        c2 = ws1.cell(row=r, column=2, value=v)
        c2.border = thin

    # Row 5: counts banner
    count_labels = ["Total", "HIGH", "MEDIUM", "LOW", "New Keywords", "Avg Sentiment"]
    count_values = [total, high_count, medium_count, low_count, len(new_kws), round(avg_sent, 3)]
    for col, (lbl, val) in enumerate(zip(count_labels, count_values), start=1):
        _hdr_cell(ws1, 5, col, lbl)
        c = ws1.cell(row=6, column=col, value=val)
        c.border = thin; c.alignment = Alignment(horizontal="center")

    ws1.freeze_panes = "A7"

    # Row 7: finding table headers (12 columns per spec)
    finding_cols = [
        "Section", "Page", "Para", "Sent", "Keyword", "Category",
        "Severity", "Flagged Sentence", "Context Before", "Context After",
        "Explanation", "Sentiment Score",
    ]
    for col, hdr in enumerate(finding_cols, start=1):
        _hdr_cell(ws1, 7, col, hdr)

    ws1.auto_filter.ref = f"A7:{get_column_letter(len(finding_cols))}7"

    sorted_findings = sorted(
        findings,
        key=lambda f: ({"HIGH": 0, "MEDIUM": 1, "LOW": 2}.get(f["severity"], 3),
                       -f["severity_score"])
    )

    high_fill   = PatternFill(start_color="FFCCCC", end_color="FFCCCC", fill_type="solid")
    medium_fill = PatternFill(start_color="FFF3CC", end_color="FFF3CC", fill_type="solid")
    low_fill    = PatternFill(start_color="E6FFE6", end_color="E6FFE6", fill_type="solid")

    data_row = 8
    for f in sorted_findings:
        sev = f["severity"]
        fill = high_fill if sev == "HIGH" else medium_fill if sev == "MEDIUM" else low_fill
        vals = [
            f["section"], f["page_num"], f["para_num"], f["sentence_num"],
            f["keyword"], f["category"], f["severity"],
            f["flagged_sentence"], f["context_before"], f["context_after"],
            f["explanation"], round(f["sentiment_score"], 3),
        ]
        for col, val in enumerate(vals, start=1):
            _data_cell(ws1, data_row, col, val, fill)
        data_row += 1

    # Column widths
    col_widths = [12, 6, 6, 6, 18, 14, 10, 60, 40, 40, 60, 10]
    for col, w in enumerate(col_widths, start=1):
        ws1.column_dimensions[get_column_letter(col)].width = w
    ws1.row_dimensions[7].height = 28

    # ── Sheet 2: Risk Dashboard ───────────────────────────────
    ws2 = wb.create_sheet("Risk Dashboard")

    metrics = [
        ("Total Findings", total),
        ("HIGH", high_count),
        ("MEDIUM", medium_count),
        ("LOW", low_count),
        ("New This Year", len(new_kws)),
        ("Avg Sentiment", round(avg_sent, 3)),
    ]
    for col, (lbl, val) in enumerate(metrics, start=1):
        _hdr_cell(ws2, 1, col, lbl)
        c = ws2.cell(row=2, column=col, value=val)
        c.border = thin; c.alignment = Alignment(horizontal="center")
        ws2.column_dimensions[get_column_letter(col)].width = 16

    # Category breakdown with YoY delta
    by_cat = comparison.get("by_category", {})
    for col, hdr in enumerate(["Category", "This Year", "Prior Year", "Delta", "Trajectory"], start=1):
        _hdr_cell(ws2, 4, col, hdr)
    r = 5
    for cat, cd in by_cat.items():
        ws2.cell(row=r, column=1, value=cat).border = thin
        ws2.cell(row=r, column=2, value=cd.get("current_count", 0)).border = thin
        ws2.cell(row=r, column=3, value=cd.get("previous_count", 0)).border = thin
        delta = cd.get("change", 0)
        dc = ws2.cell(row=r, column=4, value=f"+{delta}" if delta > 0 else str(delta))
        dc.border = thin
        ws2.cell(row=r, column=5, value=cd.get("trajectory", "")).border = thin
        r += 1

    # New keywords table
    r += 2
    _hdr_cell(ws2, r, 1, "New Keyword"); _hdr_cell(ws2, r, 2, "Category")
    r += 1
    kw_to_cat = {f["keyword"]: f["category"] for f in findings}
    for kw in new_kws:
        ws2.cell(row=r, column=1, value=kw).border = thin
        ws2.cell(row=r, column=2, value=kw_to_cat.get(kw, "")).border = thin
        r += 1

    # ── Sheet 3: YoY Comparison ───────────────────────────────
    ws3 = wb.create_sheet("YoY Comparison")

    overall = comparison.get("overall", {})
    ws3.cell(row=1, column=1, value="Year-Over-Year Risk Comparison").font = Font(bold=True, size=12)
    traj = overall.get("trajectory", "N/A")
    ws3.cell(row=2, column=1, value=f"Trajectory: {traj}").font = Font(bold=True)

    for col, hdr in enumerate(["Category", "Prior Year", "This Year", "Delta", "Trajectory"], start=1):
        _hdr_cell(ws3, 4, col, hdr)
    r = 5
    for cat, cd in by_cat.items():
        delta = cd.get("change", 0)
        for col, val in enumerate(
            [cat, cd.get("previous_count", 0), cd.get("current_count", 0),
             f"+{delta}" if delta > 0 else str(delta), cd.get("trajectory", "")],
            start=1
        ):
            ws3.cell(row=r, column=col, value=val).border = thin
        r += 1

    # New keywords
    r += 2
    _hdr_cell(ws3, r, 1, "New Keywords This Year")
    r += 1
    for kw in new_kws:
        ws3.cell(row=r, column=1, value=kw).border = thin; r += 1

    # Resolved keywords (appeared in prior but not current)
    current_kws = {f["keyword"] for f in findings}
    resolved = sorted(comparison.get("new_keywords_previous", []))
    if not resolved and comparison.get("by_category"):
        pass  # resolved list not always computed — leave empty

    r += 1
    _hdr_cell(ws3, r, 1, "Sentiment Trend")
    r += 1
    st = comparison.get("sentiment_trend", {})
    ws3.cell(row=r, column=1, value=f"Direction: {st.get('direction','N/A')}").font = Font(bold=True)
    r += 1
    ws3.cell(row=r, column=1, value=f"Prior Year: {st.get('previous', 0):.3f}  |  This Year: {st.get('current', 0):.3f}")

    for col in range(1, 6):
        ws3.column_dimensions[get_column_letter(col)].width = 20

    filepath = os.path.join(OUTPUT_DIR, f"{ticker}_redflag.xlsx")
    wb.save(filepath)
    print(f"[RedFlag] Excel saved: {filepath}")
    return filepath


# ─────────────────────────────────────────────────────────────
#  PDF  (WeasyPrint + Jinja2 → FPDF2 fallback)
# ─────────────────────────────────────────────────────────────

def _build_pdf_context(ticker, analysis, comparison, latest_date, exchange, previous_date):
    findings = analysis["findings"]
    summary  = analysis["summary"]
    high_count   = summary["by_severity"].get("HIGH", 0)
    medium_count = summary["by_severity"].get("MEDIUM", 0)
    low_count    = summary["by_severity"].get("LOW", 0)
    total        = summary["total"]

    if high_count > 0:
        verdict_label = "CRITICAL RISK"
        verdict_class = "verdict-red"
        verdict_line1 = f"{high_count} HIGH-severity flag(s) detected — immediate review required."
        verdict_line2 = "Do not proceed without resolving these items in due diligence."
    elif medium_count > 3:
        verdict_label = "ELEVATED RISK"
        verdict_class = "verdict-amber"
        verdict_line1 = f"{medium_count} MEDIUM-severity flags detected — detailed review recommended."
        verdict_line2 = "Each flag should be traced to source documents and management commentary."
    else:
        verdict_label = "LOW RISK"
        verdict_class = "verdict-green"
        verdict_line1 = f"Minimal findings ({total} total). Standard due diligence sufficient."
        verdict_line2 = "No high-priority items requiring escalation at this time."

    avg_sev = 0
    if total > 0:
        avg_sev = round(sum(f.get("severity_score", 50) for f in findings) / total)

    overall = comparison.get("overall", {})

    return {
        "ticker":          ticker,
        "exchange":        exchange,
        "filing_date":     str(latest_date or "N/A"),
        "previous_date":   str(previous_date or "N/A"),
        "generated_date":  datetime.now().strftime("%Y-%m-%d"),
        "verdict_label":   verdict_label,
        "verdict_class":   verdict_class,
        "verdict_line1":   verdict_line1,
        "verdict_line2":   verdict_line2,
        "total_findings":  total,
        "high_count":      high_count,
        "medium_count":    medium_count,
        "low_count":       low_count,
        "avg_severity":    avg_sev,
        "avg_sentiment":   summary.get("avg_sentiment", 0.0),
        "sections_scanned": summary.get("sections_with_flags", []),
        "high_findings":   [f for f in findings if f["severity"] == "HIGH"],
        "medium_findings": [f for f in findings if f["severity"] == "MEDIUM"],
        "low_findings":    [f for f in findings if f["severity"] == "LOW"],
        "has_previous":    bool(previous_date),
        "category_grid":   comparison.get("by_category", {}),
        "new_keywords":    comparison.get("new_keywords", []),
        "sentiment_trend": comparison.get("sentiment_trend", {}),
        "overall":         overall,
    }


def generate_pdf(ticker, analysis, comparison, latest_date, exchange, previous_date):
    ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, f"{ticker}_redflag.pdf")
    ctx = _build_pdf_context(ticker, analysis, comparison, latest_date, exchange, previous_date)

    # ── Try WeasyPrint ────────────────────────────────────────
    try:
        from weasyprint import HTML as WP_HTML

        template_dir = os.path.dirname(os.path.abspath(__file__))
        env = jinja2.Environment(
            loader=jinja2.FileSystemLoader(template_dir),
            autoescape=False,
        )
        html_str = env.get_template("pdf_template.html").render(**ctx)
        WP_HTML(string=html_str, base_url=template_dir).write_pdf(filepath)
        print(f"[RedFlag] PDF saved (WeasyPrint): {filepath}")
        return filepath

    except Exception as wp_err:
        print(f"[RedFlag] WeasyPrint unavailable ({wp_err}). Falling back to FPDF2.")

    # ── FPDF2 fallback ────────────────────────────────────────
    try:
        from fpdf import FPDF

        def _s(text):
            """ASCII-safe string for FPDF2 (strips non-latin-1 chars)."""
            return str(text).encode("latin-1", "replace").decode("latin-1")

        findings = analysis["findings"]
        summary  = analysis["summary"]

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        # Cover
        pdf.set_font("Helvetica", "B", 20)
        pdf.cell(0, 10, _s(f"{ticker}: RedFlag Risk Report"), ln=True)
        pdf.set_font("Helvetica", "", 11)
        pdf.cell(0, 7, _s(f"Source: {exchange}  |  Filing: {latest_date}  |  Generated: {datetime.now().strftime('%Y-%m-%d')}"), ln=True)
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 8, _s(ctx["verdict_label"]), ln=True)
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(0, 5, _s(ctx["verdict_line1"]))
        pdf.multi_cell(0, 5, _s(ctx["verdict_line2"]))
        pdf.ln(3)
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 7, "Finding Counts", ln=True)
        pdf.set_font("Helvetica", "", 11)
        hi = summary["by_severity"].get("HIGH", 0)
        me = summary["by_severity"].get("MEDIUM", 0)
        lo = summary["by_severity"].get("LOW", 0)
        pdf.cell(0, 6, f"Total: {summary['total']}  |  HIGH: {hi}  |  MEDIUM: {me}  |  LOW: {lo}", ln=True)

        # Findings (all)
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, "HIGH Severity Findings", ln=True)
        pdf.set_font("Helvetica", "", 9)
        for f in [x for x in findings if x["severity"] == "HIGH"]:
            pdf.multi_cell(0, 4, _s(f"[{f['section']} p.{f['page_num']}] {f['keyword']}: {f['flagged_sentence'][:120]}"))
            pdf.multi_cell(0, 4, _s(f"  -> {f['explanation'][:100]}"))
            pdf.ln(1)

        pdf.add_page()
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, "MEDIUM Severity Findings", ln=True)
        pdf.set_font("Helvetica", "", 9)
        for f in [x for x in findings if x["severity"] == "MEDIUM"]:
            pdf.multi_cell(0, 4, _s(f"[{f['section']} p.{f['page_num']}] {f['keyword']}: {f['flagged_sentence'][:120]}"))
            pdf.multi_cell(0, 4, _s(f"  -> {f['explanation'][:100]}"))
            pdf.ln(1)

        pdf.add_page()
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, "LOW Severity Findings", ln=True)
        pdf.set_font("Helvetica", "", 9)
        for f in [x for x in findings if x["severity"] == "LOW"]:
            pdf.multi_cell(0, 4, _s(f"[{f['section']} p.{f['page_num']}] {f['keyword']}: {f['flagged_sentence'][:120]}"))
            pdf.ln(1)

        # Methodology
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, "Methodology", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5, "All 18 standard 10-K items scanned with 220+ risk keywords across 6 categories. "
                              "Severity scored 0-100 based on sentiment, mitigating/amplifying language, "
                              "and keyword co-occurrence within sections.")
        pdf.ln(3)
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(0, 5, "github.com/zshqv/RedFlag  |  Not financial advice.", ln=True)

        pdf.output(filepath)
        print(f"[RedFlag] PDF saved (FPDF2 fallback): {filepath}")
        return filepath

    except Exception as e:
        print(f"[RedFlag] PDF generation failed: {e}")
        return None


# ─────────────────────────────────────────────────────────────
#  PPTX (python-pptx only, no matplotlib)
# ─────────────────────────────────────────────────────────────

def _blank_slide(prs, bg_rgb=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_rgb:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = bg_rgb
    return slide


def _textbox(slide, left, top, width, height, text, size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, word_wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = sanitize(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def _rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.width = 0
    return shape


def _sev_color(severity):
    return {"HIGH": RED, "MEDIUM": AMBER, "LOW": GREEN}.get(severity, GREY)


def generate_pptx(ticker, analysis, comparison, latest_date, exchange):
    ensure_output_dir()

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    findings = analysis["findings"]
    summary  = analysis["summary"]
    high_count   = summary["by_severity"].get("HIGH", 0)
    medium_count = summary["by_severity"].get("MEDIUM", 0)
    low_count    = summary["by_severity"].get("LOW", 0)
    total        = summary["total"]
    new_kws      = comparison.get("new_keywords", [])
    overall      = comparison.get("overall", {})
    by_cat       = summary.get("by_category", {})

    # ── Slide 1: Cover ───────────────────────────────────────
    print("[RedFlag] Building slide 1: Cover")
    s1 = _blank_slide(prs, NAVY)
    _textbox(s1, Inches(1), Inches(1.8), Inches(11.333), Inches(2),
             f"{ticker}: RedFlag Risk Report", size=44, bold=True, color=WHITE)
    _textbox(s1, Inches(1), Inches(3.6), Inches(8), Inches(0.7),
             f"{exchange}  |  {latest_date}", size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Verdict box (bottom right)
    if high_count > 0:
        vl, vc = "CRITICAL RISK", RED
    elif medium_count > 3:
        vl, vc = "ELEVATED RISK", AMBER
    else:
        vl, vc = "LOW RISK", GREEN
    vbox = _rect(s1, Inches(9), Inches(4.8), Inches(3.8), Inches(1.5), vc)
    tf = vbox.text_frame
    tf.text = sanitize(vl)
    tf.paragraphs[0].font.size = Pt(22); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1

    _textbox(s1, Inches(1), Inches(5.0), Inches(7.5), Inches(0.6),
             f"Total: {total}  |  HIGH: {high_count}  |  MEDIUM: {medium_count}  |  LOW: {low_count}",
             size=14, color=WHITE)
    _textbox(s1, Inches(10.5), Inches(6.9), Inches(2.5), Inches(0.5),
             "RedFlag", size=12, bold=True, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)

    # ── Slide 2: Risk Scorecard ──────────────────────────────
    print("[RedFlag] Building slide 2: Risk Scorecard")
    s2 = _blank_slide(prs, GREY)
    _textbox(s2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Risk Scorecard", size=28, bold=True, color=NAVY)

    metrics = [
        ("Total",       str(total),        NAVY),
        ("HIGH",        str(high_count),   RED),
        ("MEDIUM",      str(medium_count), AMBER),
        ("LOW",         str(low_count),    GREEN),
        ("New Keywords",str(len(new_kws)), NAVY),
        ("Avg Sentiment",
         str(round(summary.get("avg_sentiment", 0.0), 3)), NAVY),
    ]
    for i, (label, value, color) in enumerate(metrics):
        lp = Inches(0.5 + i * 2.05)
        box = _rect(s2, lp, Inches(1.5), Inches(1.9), Inches(3.5), WHITE, color)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = sanitize(value)
        p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        _textbox(s2, lp, Inches(5.3), Inches(1.9), Inches(0.6),
                 label, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Slide 3: Category Risk Heat Map (native pptx chart) ──
    print("[RedFlag] Building slide 3: Category Risk Heat Map")
    s3 = _blank_slide(prs, WHITE)
    _textbox(s3, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Risk Heat Map — Findings by Category", size=26, bold=True, color=NAVY)

    if by_cat:
        cd = ChartData()
        cd.categories = list(by_cat.keys())
        cd.add_series("Risk Flags", [int(v) for v in by_cat.values()])
        chart_frame = s3.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(1.0), Inches(12.333), Inches(6.0), cd
        )
        chart = chart_frame.chart
        chart.has_legend = False
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = NAVY

    # ── Slide 4 (+ 4b): All Flags by Tier ───────────────────
    print("[RedFlag] Building slide 4: All Flags by Tier")

    high_f   = [f for f in findings if f["severity"] == "HIGH"]
    medium_f = [f for f in findings if f["severity"] == "MEDIUM"]
    low_f    = [f for f in findings if f["severity"] == "LOW"]

    MAX_PER_COL = 20
    need_overflow = any(len(t) > MAX_PER_COL for t in [high_f, medium_f, low_f])

    def _flags_slide(prs, hi_slice, me_slice, lo_slice, title):
        s = _blank_slide(prs, WHITE)
        _textbox(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 title, size=24, bold=True, color=NAVY)
        for ci, (tier, tier_f, col_color) in enumerate([
            ("HIGH", hi_slice, RED), ("MEDIUM", me_slice, AMBER), ("LOW", lo_slice, GREEN)
        ]):
            lp = Inches(0.4 + ci * 4.2)
            hbox = _rect(s, lp, Inches(1.0), Inches(3.9), Inches(0.45), col_color)
            tf = hbox.text_frame
            tf.text = sanitize(f"{tier}  ({len(tier_f)})")
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = 1

            tb = s.shapes.add_textbox(lp, Inches(1.55), Inches(3.9), Inches(5.8))
            tf2 = tb.text_frame; tf2.word_wrap = True
            for f in tier_f:
                pg = tf2.add_paragraph()
                pg.text = sanitize(f"Pg{f['page_num']} — {f['keyword']} — {f['section']}")
                pg.font.size = Pt(9)
        return s

    _flags_slide(prs, high_f[:MAX_PER_COL], medium_f[:MAX_PER_COL], low_f[:MAX_PER_COL],
                 "All Flags by Tier")
    if need_overflow:
        print("[RedFlag] Building slide 4b: All Flags by Tier (continued)")
        _flags_slide(prs, high_f[MAX_PER_COL:], medium_f[MAX_PER_COL:], low_f[MAX_PER_COL:],
                     "All Flags by Tier (continued)")

    # ── Slide 5: Top 3 Flags ────────────────────────────────
    print("[RedFlag] Building slide 5: Top 3 Flags to Present")
    s5 = _blank_slide(prs, WHITE)
    _textbox(s5, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Top 3 Flags to Present", size=26, bold=True, color=NAVY)

    for idx, f in enumerate(findings[:3]):
        tp = Inches(1.1 + idx * 2.1)
        col = _sev_color(f["severity"])
        # colored left border
        _rect(s5, Inches(0.5), tp, Inches(0.08), Inches(1.9), col)
        card = _rect(s5, Inches(0.62), tp, Inches(12.2), Inches(1.9), WHITE, NAVY)
        tf = card.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = sanitize(f"{idx+1}. {f['keyword']}  —  {f['section']}  (Pg {f['page_num']})  [{f['category']}]")
        p1.font.size = Pt(11); p1.font.bold = True; p1.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = sanitize(f["flagged_sentence"][:120])
        p2.font.size = Pt(9)
        p3 = tf.add_paragraph()
        p3.text = sanitize(f["explanation"][:140])
        p3.font.size = Pt(9); p3.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ── Slide 6: New Keywords ────────────────────────────────
    print("[RedFlag] Building slide 6: New Keywords This Year")
    s6 = _blank_slide(prs, WHITE)
    _textbox(s6, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "New Keywords This Year", size=26, bold=True, color=NAVY)
    _textbox(s6, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
             "Absent in prior year's filing", size=13, color=RGBColor(0x66, 0x66, 0x66))

    kw_to_cat = {f["keyword"]: f["category"] for f in findings}
    for i, kw in enumerate(new_kws[:24]):
        col_idx = i % 4
        row_idx = i // 4
        lp = Inches(0.5 + col_idx * 3.2)
        tp = Inches(1.5  + row_idx * 1.2)
        pill = _rect(s6, lp, tp, Inches(3.0), Inches(0.7), RGBColor(0xE8, 0xE8, 0xE8))
        tf = pill.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = sanitize(kw)
        p.font.size = Pt(10); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = sanitize(kw_to_cat.get(kw, ""))
        p2.font.size = Pt(8); p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.alignment = PP_ALIGN.CENTER

    # ── Slide 7: YoY Trajectory ──────────────────────────────
    print("[RedFlag] Building slide 7: YoY Trajectory")
    s7 = _blank_slide(prs, WHITE)
    _textbox(s7, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "Year-Over-Year Risk Trajectory", size=26, bold=True, color=NAVY)

    traj = overall.get("trajectory", "STABLE")
    traj_col = RED if traj in ("DETERIORATING", "INCREASING") else GREEN if "IMPROV" in traj else GREY
    tbox = _rect(s7, Inches(2.5), Inches(1.2), Inches(8.333), Inches(1.8), traj_col)
    tf = tbox.text_frame; tf.text = sanitize(traj)
    tf.paragraphs[0].font.size = Pt(42); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1

    # Prior vs This Year comparison table
    cat_comp = comparison.get("by_category", {})
    tb = s7.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(4.0))
    tf2 = tb.text_frame; tf2.word_wrap = True
    header_p = tf2.paragraphs[0]
    header_p.text = sanitize("Category                Prior Year  →  This Year   Delta")
    header_p.font.size = Pt(10); header_p.font.bold = True; header_p.font.color.rgb = NAVY
    for cat, cd in list(cat_comp.items())[:8]:
        pg = tf2.add_paragraph()
        delta = cd.get("change", 0)
        arrow = "▲" if delta > 0 else ("▼" if delta < 0 else "━")
        pg.text = sanitize(f"{cat:<28} {cd.get('previous_count',0):>4}  {arrow}  {cd.get('current_count',0):>4}   ({'+' if delta>0 else ''}{delta})")
        pg.font.size = Pt(10)

    # ── Slide 8: Source & Methodology ───────────────────────
    print("[RedFlag] Building slide 8: Source and Methodology")
    s8 = _blank_slide(prs, NAVY)
    _textbox(s8, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "Source & Methodology", size=26, bold=True, color=WHITE)

    lines = [
        f"Filing: {ticker}  |  {exchange}  |  {latest_date}",
        "Keywords Scanned: 220+ across 6 risk categories",
        "Sections Analyzed: All 18 standard 10-K items",
        "Severity Scoring: Rule-based, clamped 0-100",
        "Explanation Engine: Claude API (claude-haiku-4-5-20251001) + fallback",
        "",
        "github.com/zshqv/RedFlag",
        "Not financial advice. For due diligence purposes only.",
    ]
    tb = s8.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        pg = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pg.text = sanitize(line)
        pg.font.size = Pt(14); pg.font.color.rgb = WHITE
        if line.startswith("Not financial"):
            pg.font.size = Pt(11); pg.font.color.rgb = RGBColor(0xFF, 0xCC, 0x00)
        pg.space_before = Pt(4)

    filepath = os.path.join(OUTPUT_DIR, f"{ticker}_dashboard.pptx")
    prs.save(filepath)
    print(f"[RedFlag] PPTX saved: {filepath}")
    return filepath


# ─────────────────────────────────────────────────────────────
#  MASTER
# ─────────────────────────────────────────────────────────────

def generate_reports(ticker, analysis, comparison, latest_date, previous_date, exchange):
    """Master function — generates Excel, PDF, and PPTX."""
    excel_path = generate_excel(ticker, analysis, comparison, latest_date, previous_date)
    pdf_path   = generate_pdf(ticker, analysis, comparison, latest_date, exchange, previous_date)
    pptx_path  = generate_pptx(ticker, analysis, comparison, latest_date, exchange)
    return {"excel": excel_path, "pdf": pdf_path, "pptx": pptx_path}
