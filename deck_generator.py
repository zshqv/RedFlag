# =============================================================================
# deck_generator.py — RedFlag Boardroom PowerPoint Deck Generator
# 6-slide CFO-ready deck using python-pptx
# =============================================================================

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# =============================================================================
# DESIGN SYSTEM
# =============================================================================

NAVY        = RGBColor(0x1A, 0x1A, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xE2, 0x4B, 0x4A)
AMBER       = RGBColor(0xEF, 0x9F, 0x27)
GREEN       = RGBColor(0x3B, 0x6D, 0x11)
SECONDARY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_RED   = RGBColor(0xFC, 0xEB, 0xEB)
LIGHT_AMBER = RGBColor(0xFA, 0xEE, 0xDA)
LIGHT_GREEN = RGBColor(0xEA, 0xF3, 0xDE)
SUBTEXT     = RGBColor(0x88, 0x88, 0xAA)
LIGHT_GREY  = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY    = RGBColor(0xE0, 0xE0, 0xE0)
RED_TEXT    = RGBColor(0xA3, 0x2D, 0x2D)
AMBER_TEXT  = RGBColor(0x85, 0x4F, 0x0B)

COMP_COLORS = [
    NAVY,
    RGBColor(0xEF, 0x9F, 0x27),
    RGBColor(0x16, 0xA0, 0x85),
    RGBColor(0x8E, 0x44, 0xAD),
]

CATEGORY_COLORS = {
    "Financial":   RGBColor(0x29, 0x80, 0xB9),
    "Legal":       RGBColor(0x8E, 0x44, 0xAD),
    "Operational": RGBColor(0x16, 0xA0, 0x85),
    "Regulatory":  RGBColor(0xD3, 0x54, 0x00),
}

SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)
OUTPUT_DIR = "output"
CATEGORIES = ["Financial", "Legal", "Operational", "Regulatory"]


# =============================================================================
# LOW-LEVEL SHAPE / TEXT HELPERS
# =============================================================================

def _new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _no_line(shape):
    """Remove the visible border from a shape via XML."""
    try:
        sp_pr = shape._element.spPr
        ln = sp_pr.find(qn("a:ln"))
        if ln is not None:
            sp_pr.remove(ln)
        ln = etree.SubElement(sp_pr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
    except Exception:
        pass


def _rect(slide, left, top, width, height, fill=None):
    """Add a solid rectangle with no border. fill=None → transparent."""
    shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    _no_line(shape)
    return shape


def _txb(slide, left, top, width, height):
    """Add an empty textbox and return it."""
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tb.text_frame.word_wrap = True
    return tb


def _run(p, text, size, color, bold=False, italic=False):
    """Add a formatted run to a paragraph."""
    run = p.add_run()
    run.text = _san(str(text))
    run.font.name   = "Calibri"
    run.font.size   = size
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    return run


def _first_para(tf, text, size, color, bold=False, italic=False,
                align=PP_ALIGN.LEFT, space_after=0):
    """Write to the first paragraph of a text frame, clearing any existing runs."""
    p = tf.paragraphs[0]
    for child in list(p._p):
        tag = child.tag
        if tag.endswith("}r") or tag.endswith("}br"):
            p._p.remove(child)
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    _run(p, text, size, color, bold=bold, italic=italic)
    return p


def _add_para(tf, text, size, color, bold=False, italic=False,
              align=PP_ALIGN.LEFT, space_after=0, space_before=0):
    """Append a new paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    if space_before:
        p.space_before = Pt(space_before)
    _run(p, text, size, color, bold=bold, italic=italic)
    return p


def _shape_label(shape, text, size, color, bold=False, align=PP_ALIGN.CENTER):
    """Write centered text into a shape's text frame (middle-anchored)."""
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    _first_para(tf, text, size, color, bold=bold, align=align)


def _san(text, max_len=None):
    """Strip non-ASCII and optionally truncate."""
    if not isinstance(text, str):
        text = str(text)
    text = text.encode("ascii", "replace").decode("ascii").replace("?", " ").strip()
    if max_len and len(text) > max_len:
        text = text[:max_len - 3] + "..."
    return text


def _severity_rgb(score):
    if score >= 70:
        return RED
    elif score >= 40:
        return AMBER
    return GREEN


def _avg_sev(findings):
    if not findings:
        return 0
    return round(sum(f.get("severity_score", 50) for f in findings) / len(findings))


def _top_cat(summary):
    if not summary:
        return "Regulatory", 0
    c = max(summary, key=summary.get)
    return c, summary[c]


# =============================================================================
# VERDICT SENTENCE GENERATORS
# =============================================================================

def _verdict_single(ticker, analysis, comparison):
    findings = analysis.get("findings", [])
    summary  = analysis.get("summary", {})
    total    = len(findings)
    top_c, top_n = _top_cat(summary)
    trend  = _san(comparison["sentiment_trend"].get("sentiment_trend", ""))
    traj   = _san(comparison["overall"].get("trajectory", "STABLE"))
    cur    = comparison["sentiment_trend"].get("current_avg_sentiment", 0)

    s1 = (f"{ticker}'s 10-K contains {total} risk findings. "
          f"{top_c} risks represent the largest concentration with {top_n} flagged instances.")

    if findings:
        top_f = findings[0]
        kw  = _san(top_f.get("keyword", "N/A"), 50)
        sec = _san(top_f.get("section", ""))
        sev = top_f.get("severity_score", 0)
        s2 = (f"The highest-severity finding (score {sev}) is '{kw}' in the {sec}, "
              f"Para {top_f.get('para_num',0)}, Sentence {top_f.get('sentence_num',0)}.")
    else:
        s2 = "No high-severity findings were identified in this filing."

    s3 = (f"Overall language: {trend}. "
          f"Avg sentiment {cur:.3f}. Trajectory: {traj}.")
    return [s1, s2, s3]


def _verdict_compare(tickers, results_list):
    if not results_list:
        return ["No data.", "No data.", "No data."]

    by_risk = sorted(results_list,
                     key=lambda r: _avg_sev(r["analysis"].get("findings", [])),
                     reverse=True)
    worst   = by_risk[0]
    w_sev   = _avg_sev(worst["analysis"].get("findings", []))
    w_total = len(worst["analysis"].get("findings", []))
    s1 = (f"Among {len(results_list)} peers, {worst['ticker']} carries the highest risk "
          f"profile: {w_total} total findings, avg severity {w_sev}.")

    # Most differentiated category
    best_spread, best_cat, max_t, min_t = 0, "Financial", "N/A", "N/A"
    for cat in CATEGORIES:
        vals = [(r["ticker"], r["analysis"].get("summary", {}).get(cat, 0)) for r in results_list]
        vs = [v for _, v in vals]
        spread = max(vs) - min(vs)
        if spread > best_spread:
            best_spread = spread
            best_cat = cat
            max_t = max(vals, key=lambda x: x[1])[0]
            min_t = min(vals, key=lambda x: x[1])[0]
    s2 = (f"{best_cat} risk shows the greatest divergence across peers: "
          f"{max_t} leads vs {min_t} (spread: {best_spread} findings).")

    # Most stable language
    sentiments = sorted(
        [(r["ticker"], r["comparison"]["sentiment_trend"].get("current_avg_sentiment", 0),
          len(r["comparison"].get("new_keywords", []))) for r in results_list],
        key=lambda x: x[1], reverse=True
    )
    stable = sentiments[0]
    s3 = (f"{stable[0]} shows the most stable disclosure language: "
          f"avg sentiment {stable[1]:.3f}, {stable[2]} new keyword(s) this cycle.")
    return [s1, s2, s3]


# =============================================================================
# SLIDE 1 — COVER
# =============================================================================

def _slide_cover(prs, title_line, subtitle_line, date_line):
    slide = _blank_slide(prs)

    # Full-bleed navy background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    # Top-left brand tag
    tb = _txb(slide, Inches(0.4), Inches(0.25), Inches(6), Inches(0.33))
    _first_para(tb.text_frame, "REDFLAG RISK INTELLIGENCE", Pt(10), SUBTEXT, bold=True)

    # Center: company / tickers
    tb = _txb(slide, Inches(0.5), Inches(2.55), Inches(12.33), Inches(1.3))
    font_sz = Pt(32) if len(title_line) <= 20 else Pt(24)
    _first_para(tb.text_frame, title_line, font_sz, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle below title
    tb = _txb(slide, Inches(0.5), Inches(3.95), Inches(12.33), Inches(0.5))
    _first_para(tb.text_frame, subtitle_line, Pt(14), WHITE, align=PP_ALIGN.CENTER)

    # Bottom-left: dates
    tb = _txb(slide, Inches(0.4), Inches(7.02), Inches(7), Inches(0.33))
    _first_para(tb.text_frame, date_line, Pt(9), SUBTEXT)

    # Bottom-right: source
    tb = _txb(slide, Inches(6.8), Inches(7.02), Inches(6.1), Inches(0.33))
    _first_para(tb.text_frame,
                "SEC EDGAR  ·  github.com/zshqv/RedFlag",
                Pt(9), SUBTEXT, align=PP_ALIGN.RIGHT)


# =============================================================================
# SLIDE 2 — EXECUTIVE VERDICT
# =============================================================================

def _metric_card(slide, left, top, width, height, label, value, sub_label):
    """Draw one metric card: light-grey bg, small label, big number, small sub."""
    _rect(slide, left, top, width, height, LIGHT_GREY)
    pad  = Inches(0.15)
    iw   = width - Inches(0.3)
    # Label
    tb = _txb(slide, left + pad, top + Pt(6), iw, Inches(0.25))
    _first_para(tb.text_frame, _san(str(label)), Pt(9), SECONDARY)
    # Value
    tb = _txb(slide, left + pad, top + Inches(0.32), iw, Inches(0.55))
    _first_para(tb.text_frame, _san(str(value)), Pt(20), NAVY, bold=True)
    # Sub-label
    tb = _txb(slide, left + pad, top + Inches(0.92), iw, Inches(0.27))
    _first_para(tb.text_frame, _san(str(sub_label)), Pt(9), SECONDARY)


def _slide_verdict(prs, sentences, metrics):
    """
    metrics = list of (label, value, sub_label) — 4 items
    sentences = list of 3 strings
    """
    slide = _blank_slide(prs)

    col_l_w = Inches(7.8)

    # Left column header bar
    _rect(slide, 0, 0, col_l_w, Inches(0.52), NAVY)
    tb = _txb(slide, Inches(0.3), Inches(0.1), col_l_w - Inches(0.4), Inches(0.38))
    _first_para(tb.text_frame, "EXECUTIVE SUMMARY", Pt(11), WHITE, bold=True)

    # Verdict sentences
    tb = _txb(slide, Inches(0.4), Inches(0.68), col_l_w - Inches(0.55), Inches(6.55))
    tf = tb.text_frame
    for i, sent in enumerate(sentences):
        if i == 0:
            _first_para(tf, sent, Pt(13), NAVY, space_after=10)
        else:
            _add_para(tf, sent, Pt(13), NAVY, space_after=10, space_before=4)

    # Vertical divider
    _rect(slide, Inches(8.05), Inches(0.3), Inches(0.012), Inches(6.9), MID_GREY)

    # Right column: 4 metric cards
    card_h   = Inches(1.33)
    card_gap = Inches(0.1)
    card_l   = Inches(8.22)
    card_w   = SLIDE_W - card_l - Inches(0.3)
    for i, (lbl, val, sub) in enumerate(metrics):
        card_t = Inches(0.35) + i * (card_h + card_gap)
        _metric_card(slide, card_l, card_t, card_w, card_h, lbl, val, sub)


# =============================================================================
# SLIDE 3 — RISK SCORECARD TABLE
# =============================================================================

def _style_cell(cell, fill, text_color, text, size=Pt(11), bold=False, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    for child in list(p._p):
        if child.tag.endswith("}r") or child.tag.endswith("}br"):
            p._p.remove(child)
    p.alignment = align
    _run(p, str(text), size, text_color, bold=bold)


def _cond_cell(cell, value, col_vals, size=Pt(11)):
    """Apply red/amber/green conditional fill based on position in column."""
    mn, mx = min(col_vals), max(col_vals)
    if mx == mn:
        fill, tc = LIGHT_AMBER, AMBER_TEXT
    elif value == mx:
        fill, tc = LIGHT_RED,   RED_TEXT
    elif value == mn:
        fill, tc = LIGHT_GREEN, GREEN
    else:
        fill, tc = LIGHT_AMBER, AMBER_TEXT
    _style_cell(cell, fill, tc, value, size)


def _slide_scorecard_single(prs, comparison):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(10), Inches(0.4))
    _first_para(tb.text_frame, "RISK SCORECARD", Pt(12), WHITE, bold=True)

    tbl_l = Inches(0.67)
    tbl_t = Inches(0.75)
    tbl_w = Inches(12.0)
    tbl_h = Inches(6.35)

    tbl_shape = slide.shapes.add_table(5, 5, int(tbl_l), int(tbl_t), int(tbl_w), int(tbl_h))
    tbl = tbl_shape.table
    for i, w in enumerate([Inches(2.4)] * 5):
        tbl.columns[i].width = int(w)

    for c, h in enumerate(["Category", "Current Year", "Prior Year", "Change", "Trend"]):
        _style_cell(tbl.cell(0, c), NAVY, WHITE, h, Pt(10), bold=True)

    by_cat = comparison.get("by_category", {})
    all_curr = [by_cat.get(c, {}).get("current_count", 0) for c in CATEGORIES]
    all_prev = [by_cat.get(c, {}).get("previous_count", 0) for c in CATEGORIES]
    all_chg  = [by_cat.get(c, {}).get("change", 0)        for c in CATEGORIES]

    for r, cat in enumerate(CATEGORIES, 1):
        d = by_cat.get(cat, {"current_count": 0, "previous_count": 0,
                              "change": 0, "trajectory": "STABLE"})
        _style_cell(tbl.cell(r, 0), LIGHT_GREY, NAVY, cat, Pt(10), bold=True,
                    align=PP_ALIGN.LEFT)
        _cond_cell(tbl.cell(r, 1), d.get("current_count", 0),  all_curr)
        _cond_cell(tbl.cell(r, 2), d.get("previous_count", 0), all_prev)
        _cond_cell(tbl.cell(r, 3), d.get("change", 0),         all_chg)
        _style_cell(tbl.cell(r, 4), LIGHT_GREY, SECONDARY,
                    _san(d.get("trajectory", "STABLE"), 20), Pt(10))


def _slide_scorecard_compare(prs, results_list):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(10), Inches(0.4))
    _first_para(tb.text_frame, "RISK SCORECARD", Pt(12), WHITE, bold=True)

    n = len(results_list)
    tbl_l = Inches(0.67)
    tbl_t = Inches(0.75)
    tbl_w = Inches(12.0)
    tbl_h = min(Inches(6.35), Inches(0.7) + n * Inches(1.1))

    tbl_shape = slide.shapes.add_table(n + 1, 7, int(tbl_l), int(tbl_t),
                                       int(tbl_w), int(tbl_h))
    tbl = tbl_shape.table
    for i, w in enumerate([Inches(1.5), Inches(1.7), Inches(1.7),
                            Inches(1.7), Inches(1.7), Inches(2.0), Inches(1.7)]):
        tbl.columns[i].width = int(w)

    hdrs = ["Company", "Financial", "Regulatory", "Operational", "Legal", "Avg Severity", "Total"]
    for c, h in enumerate(hdrs):
        _style_cell(tbl.cell(0, c), NAVY, WHITE, h, Pt(10), bold=True)

    data_rows = []
    for r in results_list:
        s = r["analysis"].get("summary", {})
        f = r["analysis"].get("findings", [])
        data_rows.append([r["ticker"],
                          s.get("Financial", 0), s.get("Regulatory", 0),
                          s.get("Operational", 0), s.get("Legal", 0),
                          _avg_sev(f), len(f)])

    for ri, dr in enumerate(data_rows, 1):
        _style_cell(tbl.cell(ri, 0), LIGHT_GREY, NAVY, dr[0], Pt(10), bold=True)
        for ci in range(1, 7):
            col_vals = [d[ci] for d in data_rows]
            _cond_cell(tbl.cell(ri, ci), dr[ci], col_vals)


# =============================================================================
# SLIDE 4 — CATEGORY BREAKDOWN
# =============================================================================

def _hbar(slide, left, top, width, bar_h, label, count, max_count, color):
    """Draw one labeled horizontal bar."""
    lbl_w   = Inches(1.8)
    cnt_w   = Inches(0.55)
    track_w = width - lbl_w - cnt_w - Inches(0.1)

    tb = _txb(slide, left, top, lbl_w, bar_h)
    try:
        from pptx.enum.text import MSO_ANCHOR
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    _first_para(tb.text_frame, _san(label, 20), Pt(9), SECONDARY)

    _rect(slide, left + lbl_w, top, track_w, bar_h, LIGHT_GREY)
    if max_count > 0 and count > 0:
        fw = int(track_w * min(count, max_count) / max_count)
        if fw > 0:
            _rect(slide, left + lbl_w, top, fw, bar_h, color)

    tb = _txb(slide, left + lbl_w + track_w + Inches(0.05), top, cnt_w, bar_h)
    try:
        from pptx.enum.text import MSO_ANCHOR
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    _first_para(tb.text_frame, str(count), Pt(10), NAVY, bold=True)


def _slide_categories_single(prs, analysis, comparison):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(10), Inches(0.4))
    _first_para(tb.text_frame, "CATEGORY BREAKDOWN", Pt(12), WHITE, bold=True)

    summary = analysis.get("summary", {})
    by_cat  = comparison.get("by_category", {})
    counts  = {cat: summary.get(cat, 0) for cat in CATEGORIES}
    max_cnt = max(counts.values()) if any(counts.values()) else 1

    def _bar_color(cat):
        traj = by_cat.get(cat, {}).get("trajectory", "STABLE")
        if "DETERIORATING" in traj or "INCREASING" in traj:
            return RED
        elif "IMPROVING" in traj:
            return GREEN
        return RGBColor(0x29, 0x80, 0xB9)

    content_top = Inches(0.65)
    avail_h     = SLIDE_H - content_top - Inches(0.2)
    cell_w      = (SLIDE_W - Inches(0.6)) / 2
    cell_h      = avail_h / 2
    margin      = Inches(0.35)
    bar_h       = Inches(0.42)

    positions = [
        (Inches(0.3),          content_top),
        (Inches(0.3) + cell_w, content_top),
        (Inches(0.3),          content_top + cell_h),
        (Inches(0.3) + cell_w, content_top + cell_h),
    ]

    for idx, cat in enumerate(CATEGORIES):
        cx, cy = positions[idx]

        # Category name
        tb = _txb(slide, cx + margin, cy + Inches(0.08), cell_w - margin * 2, Inches(0.33))
        _first_para(tb.text_frame, cat.upper(), Pt(10), NAVY, bold=True)

        # Bar
        count = counts.get(cat, 0)
        _hbar(slide, cx + margin, cy + Inches(0.48),
              cell_w - margin * 2, bar_h, "", count, max_cnt, _bar_color(cat))

        # Prominent count
        tb = _txb(slide, cx + margin, cy + Inches(1.05), cell_w - margin * 2, Inches(0.55))
        _first_para(tb.text_frame, f"{count} findings", Pt(17), _bar_color(cat), bold=True)


def _slide_categories_compare(prs, results_list):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(10), Inches(0.4))
    _first_para(tb.text_frame, "CATEGORY BREAKDOWN", Pt(12), WHITE, bold=True)

    # Skip Operational if all zeros
    cats_show = []
    for cat in ["Financial", "Legal", "Regulatory", "Operational"]:
        vals = [r["analysis"].get("summary", {}).get(cat, 0) for r in results_list]
        if any(v > 0 for v in vals):
            cats_show.append(cat)
    if not cats_show:
        cats_show = ["Financial", "Legal", "Regulatory"]

    n_cats    = len(cats_show)
    n_comps   = len(results_list)
    cont_top  = Inches(0.65)
    avail_h   = SLIDE_H - cont_top - Inches(0.15)
    block_h   = avail_h / n_cats
    bar_h     = max(Inches(0.28), (block_h - Inches(0.38)) / n_comps)
    chart_w   = SLIDE_W - Inches(0.8)

    for ci, cat in enumerate(cats_show):
        blk_top = cont_top + ci * block_h

        # Category label bar
        cat_bg = _rect(slide, Inches(0.4), blk_top, chart_w, Inches(0.3),
                       RGBColor(0xEC, 0xEC, 0xF4))
        tb = _txb(slide, Inches(0.55), blk_top + Inches(0.02), chart_w - Inches(0.3), Inches(0.26))
        _first_para(tb.text_frame, cat.upper(), Pt(9), NAVY, bold=True)

        vals    = [r["analysis"].get("summary", {}).get(cat, 0) for r in results_list]
        max_val = max(vals) if vals else 1

        for ri, r in enumerate(results_list):
            row_t = blk_top + Inches(0.33) + ri * bar_h
            count = r["analysis"].get("summary", {}).get(cat, 0)
            color = COMP_COLORS[ri % len(COMP_COLORS)]
            _hbar(slide, Inches(0.4), row_t, chart_w, bar_h - Inches(0.02),
                  r["ticker"], count, max_val, color)


# =============================================================================
# SLIDE 5 — TOP 5 FINDINGS
# =============================================================================

def _slide_top5(prs, findings_with_ticker):
    """
    findings_with_ticker = list of (ticker_or_None, finding_dict)
    sorted by severity desc, max 5 items
    """
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(12), Inches(0.4))
    _first_para(tb.text_frame, "HIGHEST SEVERITY FINDINGS", Pt(12), WHITE, bold=True)

    if not findings_with_ticker:
        tb = _txb(slide, Inches(0.4), Inches(2), Inches(12), Inches(1))
        _first_para(tb.text_frame, "No findings to display.", Pt(12), SECONDARY, italic=True)
        return

    n_items   = min(5, len(findings_with_ticker))
    row_h     = (SLIDE_H - Inches(0.63)) / n_items
    bar_w     = Inches(0.07)
    cont_left = Inches(0.35)
    cont_w    = SLIDE_W - Inches(0.65)

    for i, (comp_tk, f) in enumerate(findings_with_ticker[:5]):
        row_top = Inches(0.63) + i * row_h

        # Severity bar
        sev  = f.get("severity_score", 50)
        sev_c = _severity_rgb(sev)
        _rect(slide, cont_left, row_top + Inches(0.05),
              bar_w, row_h - Inches(0.12), sev_c)

        tx_left = cont_left + bar_w + Inches(0.1)
        tx_w    = cont_w - bar_w - Inches(0.12)
        kw      = _san(f.get("keyword",  ""), 45)
        cat     = _san(f.get("category", ""), 20)
        sec     = _san(f.get("section",  ""), 25)
        lbl     = _san(f.get("sentiment_label", ""), 20)
        para    = f.get("para_num",     0)
        sent    = f.get("sentence_num", 0)
        sentence = _san(f.get("sentence", ""), 200)

        # Row 1: company badge (compare) + keyword + badges
        line1_t = row_top + Inches(0.06)
        x_cursor = tx_left

        if comp_tk:
            badge = _rect(slide, x_cursor, line1_t, Inches(0.65), Inches(0.28), NAVY)
            _shape_label(badge, comp_tk, Pt(8), WHITE, bold=True)
            x_cursor += Inches(0.7)

        tb = _txb(slide, x_cursor, line1_t, Inches(3.4), Inches(0.3))
        _first_para(tb.text_frame, kw, Pt(11), NAVY, bold=True)
        x_cursor += Inches(3.5)

        cat_c = CATEGORY_COLORS.get(cat, NAVY)
        cat_badge = _rect(slide, x_cursor, line1_t, Inches(1.3), Inches(0.27), cat_c)
        _shape_label(cat_badge, cat, Pt(8), WHITE, bold=True)
        x_cursor += Inches(1.35)

        sev_badge = _rect(slide, x_cursor, line1_t, Inches(0.75), Inches(0.27), sev_c)
        _shape_label(sev_badge, str(sev), Pt(9), WHITE, bold=True)

        # Row 2: location
        loc = f"{sec}  |  Para {para}  |  Sentence {sent}"
        tb = _txb(slide, tx_left, row_top + Inches(0.38), tx_w, Inches(0.25))
        _first_para(tb.text_frame, loc, Pt(9), SECONDARY)

        # Row 3: sentence
        tb = _txb(slide, tx_left + Inches(0.15), row_top + Inches(0.63),
                  tx_w - Inches(0.2), row_h - Inches(0.72))
        _first_para(tb.text_frame, sentence, Pt(9), SECONDARY, italic=True)

        # Separator
        if i < n_items - 1:
            _rect(slide, cont_left, row_top + row_h - Inches(0.008),
                  cont_w, Inches(0.008), MID_GREY)


# =============================================================================
# SLIDE 6 — NEW KEYWORDS
# =============================================================================

def _draw_chips(slide, keywords, kw_lookup, start_left, start_top, avail_w, avail_h):
    """Render a 4-column grid of keyword chips. Returns rows drawn."""
    chip_w  = (avail_w - Inches(0.3)) / 4
    chip_h  = Inches(0.44)
    gap_y   = Inches(0.1)
    cols    = 4
    max_rows = int(avail_h / (chip_h + gap_y))

    cx, cy = 0, 0
    for kw in keywords:
        if cy >= max_rows:
            break
        cat, sev = kw_lookup.get(kw, ("Unknown", 50))
        cat_c = CATEGORY_COLORS.get(cat, NAVY)
        cat_abbr = cat[:3].upper()

        cl = start_left + cx * (chip_w + Inches(0.1))
        ct = start_top  + cy * (chip_h + gap_y)

        # Chip background
        _rect(slide, cl, ct, chip_w, chip_h, LIGHT_GREY)

        # Keyword text
        tb = _txb(slide, cl + Inches(0.1), ct + Inches(0.06),
                  chip_w - Inches(0.85), chip_h - Inches(0.1))
        try:
            from pptx.enum.text import MSO_ANCHOR
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        _first_para(tb.text_frame, _san(kw, 25), Pt(9), NAVY)

        # Category badge (right of chip)
        bw = Inches(0.65)
        cat_badge = _rect(slide, cl + chip_w - bw - Inches(0.04),
                          ct + Inches(0.07), bw, chip_h - Inches(0.14), cat_c)
        _shape_label(cat_badge, cat_abbr, Pt(7), WHITE, bold=True)

        cx += 1
        if cx >= cols:
            cx = 0
            cy += 1
    return cy + (1 if cx > 0 else 0)


def _slide_new_keywords_single(prs, comparison, findings):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(12.5), Inches(0.4))
    _first_para(tb.text_frame,
                "NEW KEYWORDS THIS CYCLE  |  NOT PRESENT IN PRIOR YEAR FILING",
                Pt(10), WHITE, bold=True)

    kw_lookup = {f.get("keyword", ""): (f.get("category", "Unknown"),
                                         f.get("severity_score", 50))
                 for f in findings}
    new_kws = comparison.get("new_keywords", [])

    cont_top  = Inches(0.65)
    cont_left = Inches(0.4)
    cont_w    = SLIDE_W - Inches(0.8)
    cont_h    = SLIDE_H - cont_top - Inches(0.55)

    if new_kws:
        _draw_chips(slide, new_kws, kw_lookup, cont_left, cont_top, cont_w, cont_h)
    else:
        tb = _txb(slide, cont_left, cont_top + Inches(0.3), cont_w, Inches(0.45))
        _first_para(tb.text_frame,
                    "No new keywords detected -- all terms also present in prior year filing.",
                    Pt(11), SECONDARY, italic=True)

    # Footer note
    footer = ("These keywords appeared in the current 10-K but were absent from the prior "
              "year filing. Their introduction may signal evolving risk disclosures.")
    tb = _txb(slide, cont_left, SLIDE_H - Inches(0.48), cont_w, Inches(0.38))
    _first_para(tb.text_frame, footer, Pt(9), SECONDARY, italic=True)


def _slide_new_keywords_compare(prs, results_list):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    tb = _txb(slide, Inches(0.4), Inches(0.09), Inches(12.5), Inches(0.4))
    _first_para(tb.text_frame,
                "NEW KEYWORDS THIS CYCLE  |  NOT PRESENT IN PRIOR YEAR FILING",
                Pt(10), WHITE, bold=True)

    n         = len(results_list)
    cont_top  = Inches(0.63)
    cont_left = Inches(0.4)
    cont_w    = SLIDE_W - Inches(0.8)
    cont_h    = SLIDE_H - cont_top - Inches(0.55)
    block_h   = cont_h / max(n, 1)

    for ri, r in enumerate(results_list):
        blk_top = cont_top + ri * block_h
        kw_lookup = {f.get("keyword", ""): (f.get("category", "Unknown"),
                                             f.get("severity_score", 50))
                     for f in r["analysis"].get("findings", [])}
        new_kws = r["comparison"].get("new_keywords", [])

        # Company subheader
        tb = _txb(slide, cont_left, blk_top, Inches(2), Inches(0.33))
        _first_para(tb.text_frame, r["ticker"], Pt(11), NAVY, bold=True)

        _draw_chips(slide, new_kws, kw_lookup,
                    cont_left + Inches(2.1), blk_top,
                    cont_w - Inches(2.1), block_h - Inches(0.05))

    footer = ("These keywords appeared in the current 10-K but were absent from the prior "
              "year filing. Their introduction may signal evolving risk disclosures.")
    tb = _txb(slide, cont_left, SLIDE_H - Inches(0.48), cont_w, Inches(0.38))
    _first_para(tb.text_frame, footer, Pt(9), SECONDARY, italic=True)


# =============================================================================
# PUBLIC INTERFACE
# =============================================================================

def _ensure_out():
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)


def generate_deck(ticker, analysis, comparison, latest_date, exchange="NYSE/NASDAQ"):
    """Build a 6-slide boardroom deck for a single ticker. Returns file path."""
    _ensure_out()
    print(f"[RedFlag] Building PowerPoint deck for {ticker}...")

    prs = _new_prs()

    # Slide 1: Cover
    date_line = (f"Filing: {_san(str(latest_date))}   |   "
                 f"Generated: {datetime.now().strftime('%B %d, %Y')}")
    _slide_cover(prs, ticker, "10-K Risk Analysis Report", date_line)

    # Slide 2: Verdict
    sentences = _verdict_single(ticker, analysis, comparison)
    findings  = analysis.get("findings", [])
    metrics   = [
        ("Total Findings",     len(findings),                    "flagged sentences"),
        ("Avg Severity Score", _avg_sev(findings),               "0 low  →  100 high"),
        ("New Keywords",       len(comparison.get("new_keywords", [])), "vs prior year filing"),
        ("Risk Trajectory",    _san(comparison["overall"].get("trajectory","STABLE"), 14),
                               "year-over-year"),
    ]
    _slide_verdict(prs, sentences, metrics)

    # Slide 3: Scorecard
    _slide_scorecard_single(prs, comparison)

    # Slide 4: Category breakdown
    _slide_categories_single(prs, analysis, comparison)

    # Slide 5: Top 5 findings
    top5 = [(None, f) for f in findings[:5]]
    _slide_top5(prs, top5)

    # Slide 6: New keywords
    _slide_new_keywords_single(prs, comparison, findings)

    filepath = os.path.join(OUTPUT_DIR, f"{ticker}_deck.pptx")
    prs.save(filepath)
    print(f"[RedFlag] Deck saved: {filepath}")
    return filepath


def generate_comparison_deck(tickers, results_list):
    """Build a 6-slide boardroom comparison deck. Returns file path."""
    _ensure_out()
    print(f"[RedFlag] Building comparison PowerPoint deck...")

    prs = _new_prs()

    # Slide 1: Cover
    title     = "   ·   ".join(tickers)
    date_line = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    _slide_cover(prs, title, "Peer Comparison  ·  10-K Risk Analysis", date_line)

    # Slide 2: Verdict
    sentences = _verdict_compare(tickers, results_list)
    all_f = [f for r in results_list for f in r["analysis"].get("findings", [])]
    worst = sorted(results_list,
                   key=lambda r: _avg_sev(r["analysis"].get("findings", [])),
                   reverse=True)[0]["ticker"] if results_list else "N/A"
    metrics = [
        ("Companies Analyzed",  len(results_list),  "peers"),
        ("Combined Findings",   len(all_f),          "total flagged sentences"),
        ("Highest Risk Peer",   worst,               "by avg severity score"),
        ("Avg Severity (All)",  _avg_sev(all_f),     "0 low  →  100 high"),
    ]
    _slide_verdict(prs, sentences, metrics)

    # Slide 3: Scorecard
    _slide_scorecard_compare(prs, results_list)

    # Slide 4: Category breakdown
    _slide_categories_compare(prs, results_list)

    # Slide 5: Top 5 findings across all companies
    combined = []
    for r in results_list:
        for f in r["analysis"].get("findings", []):
            combined.append((r["ticker"], f))
    combined.sort(key=lambda x: x[1].get("severity_score", 0), reverse=True)
    _slide_top5(prs, combined[:5])

    # Slide 6: New keywords per company
    _slide_new_keywords_compare(prs, results_list)

    fname    = f"COMPARISON_{'_'.join(tickers)}_deck.pptx"
    filepath = os.path.join(OUTPUT_DIR, fname)
    prs.save(filepath)
    print(f"[RedFlag] Comparison deck saved: {filepath}")
    return filepath
